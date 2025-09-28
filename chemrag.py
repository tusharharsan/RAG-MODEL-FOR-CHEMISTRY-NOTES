import streamlit as st
import os
import time
from pptx import Presentation
import pytesseract
from PIL import Image
import io
import google.generativeai as genai
import json
import PyPDF2
from sentence_transformers import SentenceTransformer
import re
from astrapy.db import AstraDB
import wikipedia
import requests
from bs4 import BeautifulSoup
from urllib.parse import quote
from dotenv import load_dotenv
import logging

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Load environment variables (still load from .env as a fallback)
load_dotenv()

# Set page configuration
st.set_page_config(
    page_title="Chemistry Q&A RAG System",
    page_icon="🧪",
    layout="wide"
)

# Apply custom styling
st.markdown("""
<style>
    .main {
        background-color: #000000;
    }
    .stApp {
        max-width: 1200px;
        margin: 0 auto;
    }
    h1 {
        color: #1e3d59;
    }
    .stButton button {
        background-color: #1e3d59;
        color: white;
    }
    .status {
        padding: 10px;
        border-radius: 5px;
        margin-bottom: 10px;
    }
    .success {
        background-color: #d4edda;
        color: #155724;
    }
    .info {
        background-color: #d1ecf1;
        color: #0c5460;
    }
    .results {
        padding: 20px;
        border-radius: 5px;
        background-color: white;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        margin-top: 20px;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state variables
if 'initialized' not in st.session_state:
    st.session_state.initialized = False
    st.session_state.api_key_set = False
    st.session_state.astra_credentials_set = False
    st.session_state.file_processed = False
    st.session_state.collection = None
    st.session_state.embedding_model = None
    st.session_state.qa_count = 0
    st.session_state.questions = []
    st.session_state.history = []
    st.session_state.document_text = None
    # New session state variables for API keys
    st.session_state.gemini_api_key = os.getenv("GEMINI_API_KEY", "")
    st.session_state.astra_db_token = os.getenv("ASTRA_DB_TOKEN", "")
    st.session_state.astra_db_api_endpoint = os.getenv("ASTRA_DB_API_ENDPOINT", "")
    st.session_state.astra_db_keyspace = os.getenv("ASTRA_DB_KEYSPACE", "")

# Astra DB functions
def get_or_create_user_collection(
    token: str,
    api_endpoint: str,
    user_id: str,
    environment: str = None,
    keyspace: str = None,
    dimension: int = 768
):
    """
    Gets an existing collection for a user or creates one if it doesn't exist.
    """
    try:
        # Create the collection name with user ID prefix/suffix
        collection_name = f"user_{user_id}_collection"
        
        # Ensure API endpoint has protocol
        if api_endpoint and not (api_endpoint.startswith('http://') or api_endpoint.startswith('https://')):
            api_endpoint = 'https://' + api_endpoint
            logger.info(f"Added https:// protocol to API endpoint: {api_endpoint}")
        
        # Create the AstraDB client - if keyspace is a placeholder or empty, don't use it
        if not keyspace or keyspace.strip() in ["", "YOUR_ASTRA_DB_KEYSPACE_HERE"]:
            # Connect without specifying a keyspace
            astra_db = AstraDB(
                token=token,
                api_endpoint=api_endpoint
            )
            logger.info("Connecting to Astra DB without a specific keyspace")
        else:
            # Connect with the specified keyspace
            astra_db = AstraDB(
                token=token,
                api_endpoint=api_endpoint,
                namespace=keyspace
            )
            logger.info(f"Connecting to Astra DB with keyspace: {keyspace}")
        
        # Check if collection exists
        collections = astra_db.get_collections()
        collection_exists = any(collection["name"] == collection_name for collection in collections["status"]["collections"])
        
        if collection_exists:
            # Collection exists, return it
            logger.info(f"Using existing collection: {collection_name}")
            return astra_db.collection(collection_name)
        
        # Collection doesn't exist, create it
        logger.info(f"Creating new collection: {collection_name}")
        return astra_db.create_collection(
            collection_name=collection_name,
            dimension=dimension
        )
    except Exception as e:
        logger.error(f"Error with Astra DB collection: {str(e)}")
        raise

def initialize_models(gemini_api_key=None, astra_token=None, astra_api_endpoint=None, keyspace=None):
    """Initialize models and database connections."""
    try:
        # Configure Gemini with user-provided API key
        if not gemini_api_key:
            gemini_api_key = st.session_state.gemini_api_key
            
        if not gemini_api_key:
            st.error("Please enter your Gemini API key.")
            return False
            
        genai.configure(api_key=gemini_api_key)
        st.session_state.model = genai.GenerativeModel('gemini-2.0-flash')
        logger.info("Gemini model initialized successfully")
        
        # Initialize Embedding Model
        st.session_state.embedding_model = SentenceTransformer('all-mpnet-base-v2')
        logger.info("Embedding model initialized successfully")
        
        # Initialize Astra DB if credentials are available
        if not astra_token:
            astra_token = st.session_state.astra_db_token
            
        if not astra_api_endpoint:
            astra_api_endpoint = st.session_state.astra_db_api_endpoint
            
        if not keyspace:
            keyspace = st.session_state.astra_db_keyspace
        
        if astra_token and astra_api_endpoint:
            try:
                user_id = "chemistry_app"  # A fixed user ID for this app
                st.session_state.collection = get_or_create_user_collection(
                    token=astra_token,
                    api_endpoint=astra_api_endpoint,
                    user_id=user_id,
                    keyspace=keyspace,
                    dimension=768  # Dimension for 'all-mpnet-base-v2' model
                )
                st.session_state.astra_credentials_set = True
                logger.info("Astra DB initialized successfully")
            except Exception as e:
                logger.error(f"Error initializing Astra DB: {str(e)}")
                st.warning(f"Astra DB initialization failed, but the app will continue to work without vector storage: {str(e)}")
        
        # Set initialized flag
        st.session_state.initialized = True
        st.session_state.api_key_set = True
        return True
    except Exception as e:
        logger.error(f"Error initializing models: {str(e)}")
        st.error(f"Error initializing models: {str(e)}")
        return False

def extract_text_from_pdf(pdf_file):
    """Extract text from PDF files."""
    try:
        reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        return None

def extract_text_from_pptx(pptx_file):
    """Extract text from PowerPoint files."""
    try:
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        return None

def extract_text_from_image(image_file):
    """Extract text from images using OCR."""
    try:
        image = Image.open(image_file)
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        logger.error(f"Error extracting text from image: {str(e)}")
        return None

def process_uploaded_file(uploaded_file):
    """Process the uploaded file and extract text based on file type."""
    if uploaded_file is None:
        return None
    
    file_extension = uploaded_file.name.split(".")[-1].lower()
    
    try:
        if file_extension == "pdf":
            return extract_text_from_pdf(uploaded_file)
        elif file_extension in ["ppt", "pptx"]:
            return extract_text_from_pptx(uploaded_file)
        elif file_extension in ["jpg", "jpeg", "png"]:
            return extract_text_from_image(uploaded_file)
        elif file_extension in ["txt"]:
            return uploaded_file.getvalue().decode("utf-8")
        else:
            st.warning(f"Unsupported file type: {file_extension}")
            return None
    except Exception as e:
        st.error(f"Error processing file: {str(e)}")
        return None

def extract_chemistry_qa_from_text(text_content):
    """Extracts chemistry questions and answers from text using Gemini."""
    if not text_content or not text_content.strip():
        return {"questions": []}

    prompt = f"""
    Extract chemistry questions and answers from the following text. 
    Format your response as a JSON object with the following structure:
    {{
        "questions": [
            {{
                "question_text": "Full text of the question",
                "options": {{
                    "A": "Text of option A",
                    "B": "Text of option B",
                    "C": "Text of option C",
                    "D": "Text of option D"
                }},
                "correct_answer": "Letter of the correct option (A, B, C, or D)",
                "explanation": "Explanation of the answer"
            }},
            ...
        ]
    }}
    
    If there are no clear questions and answers, return an empty questions array.
    If options are not lettered, assign them letters in order.
    
    Text: {text_content[:5000]}  # Limit text to prevent token limit issues
    """

    try:
        response = st.session_state.model.generate_content(prompt)
        response_text = response.text
        
        # Try to parse as JSON
        try:
            # Find JSON content within response (in case there's extra text)
            json_match = re.search(r'```json\s*([\s\S]*?)\s*```', response_text)
            if json_match:
                json_str = json_match.group(1)
            else:
                json_str = response_text
                
            # Clean up any markdown formatting
            json_str = re.sub(r'```.*?```', '', json_str, flags=re.DOTALL)
            
            # Parse the JSON
            result = json.loads(json_str)
            return result
        except json.JSONDecodeError:
            # If parsing fails, try to identify and fix common JSON issues
            cleaned_text = response_text.strip()
            # Remove markdown code blocks if present
            cleaned_text = re.sub(r'```json|```', '', cleaned_text)
            
            # Attempt to parse again
            try:
                result = json.loads(cleaned_text)
                return result
            except:
                logger.warning(f"Failed to parse JSON from response")
                # Return a structured empty result
                return {"questions": []}
    except Exception as e:
        logger.error(f"Error extracting Q&A: {str(e)}")
        return {"questions": []}

def store_qa_in_db(questions):
    """Stores the Q&A data in Astra DB."""
    if not st.session_state.astra_credentials_set or st.session_state.collection is None:
        st.warning("Astra DB is not configured. Questions will not be stored.")
        return False
    
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    # Clear collection by deleting all documents (if any)
    try:
        # Use delete_many to remove all documents in the collection
        st.session_state.collection.delete_many({})
        logger.info("Collection cleared successfully.")
    except Exception as e:
        logger.warning(f"Could not clear existing collection: {e}")
    
    status_text.text("Storing questions in Astra DB...")
    
    for i, q in enumerate(questions):
        # Update progress
        progress_bar.progress((i + 0.5) / len(questions))
        
        # Prepare data for embedding
        question_id = f"q{i+1}"
        
        # Create full text representation for embedding
        full_text = f"Question: {q['question_text']}\n"
        
        if "options" in q:
            options_text = ""
            for option_key, option_text in q["options"].items():
                options_text += f"Option {option_key}: {option_text}\n"
            full_text += options_text
        
        if "correct_answer" in q:
            full_text += f"Correct Answer: {q['correct_answer']}\n"
        
        if "explanation" in q:
            full_text += f"Explanation: {q['explanation']}\n"
        
        # Generate embeddings
        embedding = st.session_state.embedding_model.encode(full_text).tolist()
        
        # Store in Astra DB
        document = {
            "text": full_text,
            "page_number": q.get("page_number", 0),
            "page_type": q.get("page_type", "unknown"),
            "question_id": question_id,
            "$vector": embedding
        }
        
        try:
            st.session_state.collection.insert_one(document)
        except Exception as e:
            logger.error(f"Error inserting document into Astra DB: {str(e)}")
            
        progress_bar.progress((i + 1) / len(questions))
    
    status_text.text(f"Successfully stored {len(questions)} questions in Astra DB.")
    time.sleep(1)  # Give users time to see the message
    status_text.empty()
    progress_bar.empty()
    
    st.session_state.qa_count = len(questions)
    st.session_state.questions = questions
    st.session_state.file_processed = True
    
    return True

def query_qa_db(query, n_results=5):
    """Queries the chemistry Q&A vector database and returns relevant results."""
    if not st.session_state.astra_credentials_set or st.session_state.collection is None:
        logger.warning("Astra DB is not configured. Skipping vector search.")
        return {"documents": [], "metadatas": []}
    
    try:
        # Generate embedding for the query
        query_embedding = st.session_state.embedding_model.encode(query).tolist()
        
        # Search in Astra DB using the new client
        results = st.session_state.collection.find_many(
            filter={},
            options={
                "sort": {
                    "$vector": query_embedding
                },
                "limit": n_results
            }
        )
        
        # Convert cursor to list of documents
        documents = list(results)
        
        # Format the results to match what the generate_answer function expects
        formatted_results = {
            "documents": [[doc.get("text", "")] for doc in documents],
            "metadatas": documents
        }
        
        return formatted_results
    except Exception as e:
        logger.error(f"Error querying Astra DB: {str(e)}")
        return {"documents": [], "metadatas": []}

def search_wikipedia(query, max_results=3):
    """Search Wikipedia for relevant chemistry information."""
    try:
        search_results = wikipedia.search(f"chemistry {query}", results=max_results)
        wiki_content = []
        
        for title in search_results:
            try:
                page = wikipedia.page(title, auto_suggest=False)
                wiki_content.append({
                    "title": page.title,
                    "summary": wikipedia.summary(title, sentences=4),
                    "url": page.url
                })
            except (wikipedia.exceptions.DisambiguationError, wikipedia.exceptions.PageError) as e:
                logger.warning(f"Wikipedia error for {title}: {str(e)}")
                continue
            
        return wiki_content
    except Exception as e:
        logger.warning(f"Wikipedia search error: {str(e)}")
        return []

def search_chemistry_websites(query):
    """Search chemistry-specific websites for relevant information."""
    # Import the functions from vectordb.py instead of using astrapy.db directly
    from vectordb import (
        get_or_create_user_collection,
        initialize_models,
        process_uploaded_file,
        store_qa_in_db,
        query_qa_db,
        generate_answer,
        rag_pipeline,
        process_text_input
    )
    import streamlit as st
    import os
    import time
    import logging
    from dotenv import load_dotenv
    import google.generativeai as genai
    import wikipedia
    import requests
    from bs4 import BeautifulSoup
    from urllib.parse import quote

    # This function doesn't interact with Astra DB so it can remain unchanged
    chemistry_sites = [
        "chemguide.co.uk",
        "chem.libretexts.org",
        "masterorganicchemistry.com",
        "rsc.org",
        "acs.org"
    ]
    
    results = []
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
    }
    
    for site in chemistry_sites[:2]:  # Limit to 2 sites to speed up response
        try:
            search_url = f"https://www.google.com/search?q=site:{site}+{quote(query)}"
            response = requests.get(search_url, headers=headers, timeout=5)
            
            if response.status_code == 200:
                soup = BeautifulSoup(response.text, 'html.parser')
                
                for result in soup.find_all('div', class_='g')[:1]:
                    title_elem = result.find('h3')
                    link_elem = result.find('a')
                    snippet_elem = result.find('div', class_='VwiC3b')
                    
                    if title_elem and link_elem and snippet_elem:
                        href = link_elem.get('href', '')
                        if href.startswith('/url?q='):
                            href = href.split('/url?q=')[1].split('&')[0]
                        
                        results.append({
                            "title": title_elem.text,
                            "url": href,
                            "snippet": snippet_elem.text
                        })
            else:
                logger.warning(f"Failed to search {site}. Status code: {response.status_code}")
                
        except Exception as e:
            logger.warning(f"Error searching {site}: {str(e)}")
            continue
            
    return results

def generate_answer(query, retrieved_results=None, user_document_text=None):
    """Generate an answer using all available sources."""
    try:
        # Get external information
        wiki_results = search_wikipedia(query)
        web_results = search_chemistry_websites(query)
        
        # Prepare context from all sources
        context = "Based on the following sources:\n\n"
        
        # Add retrieved Q&A from database if available
        if retrieved_results and retrieved_results.get("documents"):
            context += "Retrieved Q&A Information:\n"
            for doc_list in retrieved_results.get("documents", []):
                if doc_list:  # Make sure it's not empty
                    context += doc_list[0] + "\n\n"
        
        # Add user document if available
        if user_document_text:
            context += "User Provided Document:\n"
            # Take the first 2000 chars to avoid token limits
            context += user_document_text[:2000] + "...\n\n"
        
        # Add Wikipedia information
        if wiki_results:
            context += "Wikipedia Information:\n"
            for wiki in wiki_results:
                context += f"- {wiki['title']}: {wiki['summary']}\n\n"
        
        # Add chemistry website information
        if web_results:
            context += "Chemistry Website Information:\n"
            for web in web_results:
                context += f"- {web['title']}: {web['snippet']}\n\n"
        
        prompt = f"""
        {context}
        
        Question: {query}
        
        Please provide:
        1. A clear, detailed answer explaining the chemistry concepts
        2. Include relevant equations or mechanisms if applicable
        3. Add practical examples where possible
        4. Mention any safety considerations if relevant
        
        Format the response in markdown with clear headings and bullet points where appropriate.
        If there are chemical equations, please format them properly.
        """
        
        try:
            response = st.session_state.model.generate_content(prompt)
            answer_text = response.text
        except Exception as e:
            logger.error(f"Error in model generation: {str(e)}")
            # Fallback with shorter context if we get token limit issues
            shortened_prompt = f"""
            Question: {query}
            
            Please provide:
            1. A clear, detailed answer explaining the chemistry concepts
            2. Include relevant equations or mechanisms if applicable
            3. Add practical examples where possible
            4. Mention any safety considerations if relevant
            """
            response = st.session_state.model.generate_content(shortened_prompt)
            answer_text = response.text
        
        # Add sources section
        answer = answer_text + "\n\n### Sources:\n"
        
        if wiki_results:
            answer += "\nWikipedia:\n"
            for wiki in wiki_results:
                answer += f"- [{wiki['title']}]({wiki['url']})\n"
        
        if web_results:
            answer += "\nChemistry Websites:\n"
            for web in web_results:
                answer += f"- [{web['title']}]({web['url']})\n"
        
        return answer
    except Exception as e:
        logger.error(f"Error generating answer: {str(e)}")
        return f"Error generating answer: {str(e)}"

def rag_pipeline(user_query, n_results=3):
    """Full RAG pipeline: Retrieval + Answer Generation."""
    # Retrieve relevant Q&A content
    with st.spinner("Searching the knowledge base..."):
        retrieved_results = query_qa_db(user_query, n_results)
    
    # Get user document if available
    user_document_text = st.session_state.get('document_text') if hasattr(st.session_state, 'document_text') else None
    
    # Generate answer based on retrieved content
    with st.spinner("Generating your answer..."):
        answer = generate_answer(user_query, retrieved_results, user_document_text)
    
    # Store in history
    st.session_state.history.append({
        "query": user_query,
        "answer": answer,
        "timestamp": time.strftime("%Y-%m-%d %H:%M:%S")
    })
    
    return answer

def process_document_for_qa(document_text):
    """Process document text and extract Q&A data."""
    status_text = st.empty()
    progress_bar = st.progress(0)
    
    status_text.text("Processing document content...")
    progress_bar.progress(0.3)
    
    # Split document into manageable chunks
    chunk_size = 5000
    overlap = 500
    chunks = []
    
    for i in range(0, len(document_text), chunk_size - overlap):
        chunk = document_text[i:i + chunk_size]
        chunks.append(chunk)
    
    progress_bar.progress(0.5)
    status_text.text(f"Extracting Q&A from document ({len(chunks)} chunks)...")
    
    all_questions = []
    
    for i, chunk in enumerate(chunks):
        status_text.text(f"Processing chunk {i+1}/{len(chunks)}...")
        qa_data = extract_chemistry_qa_from_text(chunk)
        
        if qa_data and "questions" in qa_data and qa_data["questions"]:
            for q in qa_data["questions"]:
                q["chunk_number"] = i+1
                all_questions.append(q)
        
        progress_bar.progress(0.5 + 0.5 * (i + 1) / len(chunks))
    
    status_text.text(f"Extracted {len(all_questions)} questions from document.")
    time.sleep(1)
    status_text.empty()
    progress_bar.empty()
    
    return all_questions

def main():
    st.title("🧪 Chemistry Q&A RAG System")
    st.markdown("Upload your chemistry documents and ask questions to get detailed answers!")
    
    # Sidebar
    with st.sidebar:
        st.header("System Setup")
        
        # Add a section for API credentials
        st.subheader("API Credentials")
        
        # Gemini API key input
        gemini_api_key = st.text_input(
            "Gemini API Key:", 
            value=st.session_state.gemini_api_key,
            type="password",
            help="Enter your Gemini API key to use the generative AI capabilities"
        )
        if gemini_api_key:
            st.session_state.gemini_api_key = gemini_api_key
        
        # Astra DB credentials
        with st.expander("Astra DB Credentials (Optional)"):
            astra_db_token = st.text_input(
                "Astra DB Token:", 
                value=st.session_state.astra_db_token,
                type="password",
                help="Enter your Astra DB token for vector storage"
            )
            if astra_db_token:
                st.session_state.astra_db_token = astra_db_token
                
            astra_db_api_endpoint = st.text_input(
                "Astra DB API Endpoint:", 
                value=st.session_state.astra_db_api_endpoint,
                help="Enter your Astra DB API endpoint URL"
            )
            if astra_db_api_endpoint:
                st.session_state.astra_db_api_endpoint = astra_db_api_endpoint
                
            astra_db_keyspace = st.text_input(
                "Astra DB Keyspace (Optional):", 
                value=st.session_state.astra_db_keyspace,
                help="Enter your Astra DB keyspace if using one"
            )
            if astra_db_keyspace:
                st.session_state.astra_db_keyspace = astra_db_keyspace
        
        # Initialize system button
        if not st.session_state.initialized:
            if st.button("Initialize System"):
                with st.spinner("Initializing..."):
                    # Use the initialize_models function from vectordb.py
                    if initialize_models(
                        gemini_api_key=st.session_state.gemini_api_key,
                        astra_token=st.session_state.astra_db_token,
                        astra_api_endpoint=st.session_state.astra_db_api_endpoint,
                        keyspace=st.session_state.astra_db_keyspace
                    ):
                        st.success("System initialized successfully!")
                        time.sleep(1)
                        st.rerun()
                    else:
                        st.error("Failed to initialize. Check your API credentials.")
        
        if st.session_state.initialized:
            st.success("System is ready!")
            
            # File upload
            st.subheader("Upload Chemistry Documents")
            uploaded_file = st.file_uploader("Upload PDF, PPTX, or image files:", 
                                            type=["pdf", "pptx", "jpg", "jpeg", "png", "txt"])
            
            if uploaded_file is not None:
                with st.spinner("Processing document..."):
                    document_text = process_uploaded_file(uploaded_file)
                    if document_text:
                        st.session_state.document_text = document_text
                        st.success(f"Processed {uploaded_file.name} successfully!")
                        
                        # Option to extract Q&A from document
                        if st.button("Extract Q&A from Document"):
                            with st.spinner("Extracting questions and answers..."):
                                questions = process_document_for_qa(document_text)
                                if questions:
                                    store_qa_in_db(questions)
                                    st.success(f"Successfully extracted {len(questions)} questions!")
                                    st.session_state.file_processed = True
                                else:
                                    st.warning("No questions were extracted from the document.")
                    else:
                        st.error("Failed to process the document.")
            
            # Direct text input option
            st.subheader("Or Enter Text Directly")
            text_area = st.text_area("Enter chemistry content:", height=150)
            
            if text_area and st.button("Process Text Input"):
                with st.spinner("Processing text input..."):
                    st.session_state.document_text = text_area
                    questions = process_document_for_qa(text_area)
                    if questions:
                        store_qa_in_db(questions)
                        st.success(f"Successfully extracted {len(questions)} questions!")
                        st.session_state.file_processed = True
                    else:
                        st.warning("No questions were extracted from the text.")
            
            # Display statistics
            if st.session_state.file_processed:
                st.markdown("---")
                st.subheader("Knowledge Base Stats")
                st.info(f"Questions in database: {st.session_state.qa_count}")
            
            if st.button("Clear Chat History"):
                st.session_state.history = []
                st.success("Chat history cleared!")
    
    # Main area
    if not st.session_state.initialized:
        st.info("👈 Please enter your API credentials and initialize the system using the sidebar.")
        
        # Display information about the system
        st.markdown("""
        ## About Chemistry Q&A RAG System
        
        This application helps chemistry students and professionals find answers to chemistry-related questions.
        
        Features:
        - Ask questions about any chemistry topic
        - Upload your own chemistry documents for context
        - Get answers with references to reliable sources
        - Track your question history
        
        To get started:
        1. Enter your Gemini API key in the sidebar
        2. (Optional) Add Astra DB credentials for vector search capabilities
        3. Click "Initialize System" in the sidebar
        4. Optionally upload relevant chemistry documents
        5. Ask your questions in the input field
        
        ### Required API Keys:
        - **Gemini API Key**: Necessary for all text generation and document processing
        - **Astra DB Credentials**: Optional, but recommended for better performance with vector search
        """)
    else:
        # Q&A Interface
        st.header("Ask Chemistry Questions")
        
        col1, col2 = st.columns([3, 1])
        with col1:
            # Query input
            user_query = st.text_input("Enter your chemistry question:", 
                                     placeholder="e.g., What is the mechanism of the Grignard reaction?")
        
        with col2:
            # Example questions dropdown
            example_questions = [
                "Select an example question",
                "What is the mechanism of the Grignard reaction?",
                "Explain the principles of acid-base titration",
                "How does the periodic table organize elements?",
                "What are the different types of chemical bonds?",
                "Explain the role of catalysts in chemical reactions"
            ]
            
            selected_example = st.selectbox("Or choose an example:", example_questions)
            
            if selected_example != "Select an example question" and not user_query:
                user_query = selected_example
        
        col1, col2 = st.columns([1, 3])
        with col1:
            ask_button = st.button("Ask Question")
        
        with col2:
            search_options = st.multiselect(
                "Information sources:", 
                ["Document Database", "Wikipedia", "Chemistry Websites"],
                default=["Document Database", "Wikipedia", "Chemistry Websites"]
            )
        
        if ask_button and user_query:
            with st.spinner("Researching and generating answer..."):
                answer = rag_pipeline(user_query)
                
                # Display answer
                tab1, tab2 = st.tabs(["Answer", "Sources"])
                
                with tab1:
                    st.markdown("### Your Question")
                    st.info(user_query)
                    st.markdown("### Answer")
                    main_answer = answer.split("### Sources:")[0] if "### Sources:" in answer else answer
                    st.markdown(main_answer)
                
                with tab2:
                    if "### Sources:" in answer:
                        sources = "### Sources:" + answer.split("### Sources:")[1]
                        st.markdown(sources)
                    else:
                        st.info("No external sources were used for this answer.")
        
        # Display conversation history
        if st.session_state.history:
            st.markdown("---")
            st.markdown("### Previous Questions")
            
            for i, item in enumerate(reversed(st.session_state.history[:-1])):
                with st.expander(f"Q: {item['query'][:80]}{'...' if len(item['query']) > 80 else ''} ({item.get('timestamp', 'Unknown time')})"):
                    st.markdown("#### Question")
                    st.info(item["query"])
                    st.markdown("#### Answer")
                    st.markdown(item["answer"])
        
        # Display usage statistics
        if len(st.session_state.history) > 0:
            st.sidebar.markdown("---")
            st.sidebar.subheader("Usage Statistics")
            st.sidebar.markdown(f"Questions asked: {len(st.session_state.history)}")
            st.sidebar.markdown(f"Session started: {time.strftime('%Y-%m-%d %H:%M:%S')}")

if __name__ == "__main__":
    main()