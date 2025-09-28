import streamlit as st
import os
import time
from pptx import Presentation
import pytesseract
from PIL import Image
import io
import google.generativeai as genai
import json
import PyPDF2
from sentence_transformers import SentenceTransformer
import re
from astrapy import DataAPIClient
from astrapy.info import CollectionDescriptor
from langchain_astradb import CollectionVectorServiceOptions
import hashlib
import uuid
from datetime import datetime, timedelta
import jwt  # pip install PyJWT
from dotenv import load_dotenv # Import load_dotenv

load_dotenv() # Load environment variables from .env file

# Set page configuration
st.set_page_config(
    page_title="Chemistry Q&A with User Authentication",
    page_icon="🧪",
    layout="wide"
)

# Apply custom styling
st.markdown("""
<style>
    .main {
        background-color: #000000;
    }
    .stApp {
        max-width: 1200px;
        margin: 0 auto;
    }
    h1 {
        color: #1e3d59;
    }
    .stButton button {
        background-color: #1e3d59;
        color: white;
    }
    .status {
        padding: 10px;
        border-radius: 5px;
        margin-bottom: 10px;
    }
    .success {
        background-color: #d4edda;
        color: #155724;
    }
    .info {
        background-color: #d1ecf1;
        color: #0c5460;
    }
    .results {
        padding: 20px;
        border-radius: 5px;
        background-color: white;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        margin-top: 20px;
    }
    /* New styles for auth pages */
    .auth-container {
        max-width: 400px;
        margin: 0 auto;
        padding: 20px;
        border-radius: 10px;
        background-color: #f8f9fa;
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
    }
    .auth-header {
        text-align: center;
        margin-bottom: 20px;
    }
    .user-info {
        background-color: #000000;
        padding: 10px;
        border-radius: 5px;
        margin-bottom: 20px;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state variables if they don't exist
if 'initialized' not in st.session_state:
    st.session_state.initialized = False
    st.session_state.api_key_set = False
    st.session_state.astra_credentials_set = False
    st.session_state.file_processed = False
    st.session_state.embedding_model = None
    st.session_state.history = []
    
    # Auth related session state
    st.session_state.user_logged_in = False
    st.session_state.user_id = None
    st.session_state.username = None
    st.session_state.jwt_token = None

# User Authentication Functions
def hash_password(password):
    """Create hashed password using SHA-256"""
    return hashlib.sha256(password.encode()).hexdigest()

def create_users_collection(token, api_endpoint, keyspace=None):
    """Create or get the users collection in Astra DB"""
    client = DataAPIClient(token=token, environment=None)
    database = client.get_database(api_endpoint=api_endpoint, token=token, keyspace=keyspace)
    
    # Check if users collection exists
    existing_collections = list(database.list_collection_names(keyspace=keyspace))
    if "users" in existing_collections:
        return database.get_collection("users", keyspace=keyspace)
    else:
        return database.create_collection("users", keyspace=keyspace)

def register_user(users_collection, username, password, email):
    """Register a new user"""
    # Check if user already exists
    existing_user = users_collection.find_one({"username": username})
    if existing_user:
        return False, "Username already exists"
    
    # Create user document
    user_id = str(uuid.uuid4())
    user_doc = {
        "user_id": user_id,
        "username": username,
        "password_hash": hash_password(password),
        "email": email,
        "created_at": datetime.now().isoformat(),
        "last_login": None
    }
    
    # Insert user
    users_collection.insert_one(user_doc)
    return True, user_id

def login_user(users_collection, username, password):
    """Login a user and return JWT token if successful"""
    user = users_collection.find_one({"username": username})
    
    if not user or user.get("password_hash") != hash_password(password):
        return False, "Invalid username or password"
    
    # Update last login time
    users_collection.update_one(
        {"username": username}, 
        {"$set": {"last_login": datetime.now().isoformat()}}
    )
    
    # Create JWT token
    token_data = {
        "user_id": user.get("user_id"),
        "username": username,
        "exp": datetime.utcnow() + timedelta(hours=24)  # 24 hour expiration
    }
    
    SECRET_KEY = "your-secret-key-change-this-in-production"  # Change this in production!
    jwt_token = jwt.encode(token_data, SECRET_KEY, algorithm="HS256")
    return True, {"jwt_token": jwt_token, "user_id": user.get("user_id")}

# Astra DB Collection Functions
def get_user_collection(token, api_endpoint, user_id, keyspace=None, dimension=768):
    """Get or create a collection for a specific user"""
    client = DataAPIClient(token=token, environment=None)
    database = client.get_database(api_endpoint=api_endpoint, token=token, keyspace=keyspace)

    # Create an even shorter collection name using a hash of the user_id
    short_id_hash = hashlib.md5(user_id.encode()).hexdigest()[:8] # First 8 chars of MD5 hash
    collection_name = f"user_{short_id_hash}_collection"

    print(f"Attempting to create collection with name: {collection_name}") # Debug print

    # Check if collection exists
    existing_collections = list(database.list_collection_names(keyspace=keyspace))
    if collection_name in existing_collections:
        return database.get_collection(collection_name, keyspace=keyspace)
    else:
        # Create the collection
        return database.create_collection(
            name=collection_name,
            keyspace=keyspace,
            dimension=dimension
        )

def check_collection_empty(collection):
    """Check if a collection is empty"""
    return collection.count_documents({}, upper_bound=10**9) == 0

# Functions for document processing
def extract_page_data(page, page_type):
    """Extracts text and image captions from a page (PPT or PDF)."""
    if page_type == "ppt":
        slide_data = {"slide_text": [], "image_captions": []}
        for shape in page.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_data["slide_text"].append(paragraph.text)
            elif shape.has_picture:
                try:
                    image_bytes = shape.image.blob
                    image = Image.open(io.BytesIO(image_bytes))
                    text = pytesseract.image_to_string(image)
                    slide_data["image_captions"].append(text)
                except Exception as e_image:
                    st.warning(f"Error processing image: {e_image}")
        return slide_data

    elif page_type == "pdf":
        page_text = page.extract_text()
        return {"slide_text": [page_text], "image_captions": []}  # PDFs don't have images in the same way.

    else:
        return {"slide_text": [], "image_captions": []}

def process_ppt_pages(ppt_file):
    """Processes each page of a PPT and extracts data."""
    try:
        presentation = Presentation(ppt_file)
        page_data = []
        for i, slide in enumerate(presentation.slides):
            data = extract_page_data(slide, "ppt")
            page_data.append({"page_number": i + 1, "data": data, "page_type": "ppt"})
        return page_data
    except Exception as e:
        st.error(f"Error processing PPT: {e}")
        return []

def process_pdf_pages(pdf_file):
    """Processes each page of a PDF and extracts data."""
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        page_data = []
        for i, page in enumerate(pdf_reader.pages):
            data = extract_page_data(page, "pdf")
            page_data.append({"page_number": i + 1, "data": data, "page_type": "pdf"})
        return page_data
    except Exception as e:
        st.error(f"Error processing PDF: {e}")
        return []

def extract_chemistry_qa_from_page(page_data):
    """Extracts chemistry questions and answers from a page's text using Gemini."""
    slide_text = "\n".join(page_data["slide_text"])
    image_captions = "\n".join(page_data["image_captions"])
    combined_text = slide_text + "\n" + image_captions

    if not combined_text.strip():
        return None

    prompt = f"""
    Extract chemistry questions and answers from the following text. 
    Format your response as a JSON object with the following structure:
    {{
        "questions": [
            {{
                "question_text": "Full text of the question",
                "options": {{
                    "A": "Text of option A",
                    "B": "Text of option B",
                    "C": "Text of option C",
                    "D": "Text of option D"
                }},
                "correct_answer": "Letter of the correct option (A, B, C, or D)",
                "explanation": "Explanation of the answer"
            }},
            ...
        ]
    }}
    
    If there are no clear questions and answers, return an empty questions array.
    If options are not lettered, assign them letters in order.
    
    Text: {combined_text}
    """

    try:
        response = st.session_state.model.generate_content(prompt)
        response_text = response.text
        
        # Try to parse as JSON
        try:
            # Find JSON content within response (in case there's extra text)
            json_match = re.search(r'```json\s*([\s\S]*?)\s*```', response_text)
            if json_match:
                json_str = json_match.group(1)
            else:
                json_str = response_text
                
            # Clean up any markdown formatting
            json_str = re.sub(r'```.*?```', '', json_str, flags=re.DOTALL)
            
            # Parse the JSON
            result = json.loads(json_str)
            return result
        except json.JSONDecodeError:
            # If parsing fails, try to identify and fix common JSON issues
            cleaned_text = response_text.strip()
            # Remove markdown code blocks if present
            cleaned_text = re.sub(r'```json|```', '', cleaned_text)
            
            # Attempt to parse again
            try:
                result = json.loads(cleaned_text)
                return result
            except:
                st.warning(f"Failed to parse JSON from response: {response_text}")
                # Return a structured empty result
                return {"questions": []}
    except Exception as e:
        st.error(f"Error extracting Q&A: {e}")
        return {"questions": []}

def process_uploaded_file(uploaded_file):
    """Process the uploaded file and extract Q&A data."""
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    
    # Create a progress bar
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    if file_extension == '.pptx':
        # Save the uploaded file temporarily
        with open("temp_file.pptx", "wb") as f:
            f.write(uploaded_file.getbuffer())
        status_text.text("Processing PowerPoint file...")
        page_data = process_ppt_pages("temp_file.pptx")
        os.remove("temp_file.pptx")  # Clean up
    elif file_extension == '.pdf':
        status_text.text("Processing PDF file...")
        page_data = process_pdf_pages(uploaded_file)
    else:
        st.error("Unsupported file type. Please upload a PDF or PPTX file.")
        return []
    
    all_questions = []
    total_pages = len(page_data)
    
    for i, page in enumerate(page_data):
        status_text.text(f"Extracting Q&A from page {i+1}/{total_pages}...")
        progress_bar.progress((i + 0.5) / total_pages)
        
        qa_data = extract_chemistry_qa_from_page(page["data"])
        
        if qa_data and "questions" in qa_data and qa_data["questions"]:
            for q in qa_data["questions"]:
                q["page_number"] = page["page_number"]
                q["page_type"] = page["page_type"]
                all_questions.append(q)
        
        progress_bar.progress((i + 1) / total_pages)
    
    status_text.text(f"Extracted {len(all_questions)} questions total.")
    time.sleep(1)  # Give users time to see the message
    status_text.empty()
    progress_bar.empty()
    
    return all_questions

def store_qa_in_db(questions, collection):
    """Stores the Q&A data in user's Astra DB collection."""
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    status_text.text("Storing questions in your personal database...")
    
    for i, q in enumerate(questions):
        # Update progress
        progress_bar.progress((i + 0.5) / len(questions))
        
        # Prepare data for embedding
        question_id = f"q{i+1}"
        
        # Create full text representation for embedding
        full_text = f"Question: {q['question_text']}\n"
        
        if "options" in q:
            options_text = ""
            for option_key, option_text in q["options"].items():
                options_text += f"Option {option_key}: {option_text}\n"
            full_text += options_text
        
        if "correct_answer" in q:
            full_text += f"Correct Answer: {q['correct_answer']}\n"
        
        if "explanation" in q:
            full_text += f"Explanation: {q['explanation']}\n"
        
        # Generate embeddings
        embedding = st.session_state.embedding_model.encode(full_text).tolist()
        
        # Store in Astra DB
        document = {
            "text": full_text,
            "page_number": q.get("page_number", 0),
            "page_type": q.get("page_type", "unknown"),
            "question_id": question_id,
            "$vector": embedding
        }
        
        collection.insert_one(document)
        progress_bar.progress((i + 1) / len(questions))
    
    status_text.text(f"Successfully stored {len(questions)} questions in your personal database.")
    time.sleep(1)  # Give users time to see the message
    status_text.empty()
    progress_bar.empty()
    
    st.session_state.file_processed = True

def query_qa_db(query, collection, n_results=5):
    """Queries the user's personal chemistry Q&A vector database and returns relevant results."""
    # Generate embedding for the query
    query_embedding = st.session_state.embedding_model.encode(query).tolist()
    
    # Search in Astra DB
    results = collection.find(
        vector=query_embedding,
        limit=n_results
    )
    
    # Convert cursor to list of documents
    documents = list(results)
    
    # Format the results to match what the generate_answer function expects
    formatted_results = {
        "documents": [[doc.get("text", "")] for doc in documents],
        "metadatas": documents
    }
    
    return formatted_results

def generate_answer(user_query, retrieved_results):
    """Generates an answer using Gemini based on retrieved results."""
    # Compile context from retrieved documents
    context_docs = []
    for doc_list in retrieved_results.get("documents", []):
        if doc_list:  # Make sure it's not empty
            context_docs.extend(doc_list)
    
    context = "\n\n".join(context_docs)
    
    prompt = f"""
    Based on the following chemistry question and answer information, provide a detailed answer to the user's query.
    
    Retrieved Information:
    {context}
    
    User Query: {user_query}
    
    Provide a clear, detailed answer with explanations of the chemistry concepts involved. If the answer involves chemical reactions, explain the mechanism where appropriate.
    """
    
    response = st.session_state.model.generate_content(prompt)
    return response.text

def rag_pipeline(user_query, user_collection, n_results=3):
    """Full RAG pipeline: Retrieval + Answer Generation for user's personal database."""
    # Retrieve relevant Q&A content
    with st.spinner("Searching your personal knowledge base..."):
        retrieved_results = query_qa_db(user_query, user_collection, n_results)
    
    # Generate answer based on retrieved content
    with st.spinner("Generating your answer..."):
        answer = generate_answer(user_query, retrieved_results)
    
    # Store in history
    st.session_state.history.append({
        "query": user_query,
        "answer": answer,
        "timestamp": datetime.now().isoformat()
    })
    
    return answer

# Initialize models and database
def initialize_models(keyspace=None):
    try:
        # Configure Gemini
        gemini_api_key = os.getenv("GEMINI_API_KEY")
        genai.configure(api_key=gemini_api_key)
        st.session_state.model = genai.GenerativeModel('gemini-2.0-flash')
        
        # Configure Tesseract (if needed)
        if os.name == 'nt':  # Windows
            pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract'
        
        # Initialize Embedding Model
        st.session_state.embedding_model = SentenceTransformer('all-mpnet-base-v2')
        
        # Create users collection
        astra_token = os.getenv("ASTRA_DB_TOKEN")
        astra_api_endpoint = os.getenv("ASTRA_DB_ENDPOINT")
        st.session_state.users_collection = create_users_collection(
            token=astra_token,
            api_endpoint=astra_api_endpoint,
            keyspace=keyspace
        )
        
        # Store Astra credentials in session state for later use
        st.session_state.astra_token = astra_token
        st.session_state.astra_api_endpoint = astra_api_endpoint
        if keyspace:
            st.session_state.keyspace = keyspace
        
        st.session_state.initialized = True
        st.session_state.api_key_set = True
        st.session_state.astra_credentials_set = True
        return True
    except Exception as e:
        st.error(f"Error initializing models and database: {e}")
        return False

# Authentication UI functions
def show_auth_page():
    st.header("🧪 Chemistry Q&A System")
    st.subheader("Login or Register to Access Your Personal Chemistry Knowledge Base")
    
    tab1, tab2 = st.tabs(["Login", "Register"])
    
    with tab1:
        with st.container():
            st.markdown("### Login")
            username = st.text_input("Username", key="login_username")
            password = st.text_input("Password", type="password", key="login_password")
            
            col1, col2, col3 = st.columns([1, 1, 1])
            with col2:
                login_button = st.button("Login", use_container_width=True)
            
            if login_button:
                if st.session_state.initialized and 'users_collection' in st.session_state:
                    with st.spinner("Authenticating..."):
                        success, result = login_user(st.session_state.users_collection, username, password)
                    
                    if success:
                        st.session_state.user_logged_in = True
                        st.session_state.jwt_token = result["jwt_token"]
                        st.session_state.user_id = result["user_id"]
                        st.session_state.username = username
                        st.success("Login successful! Loading your personal knowledge base...")
                        time.sleep(1)
                        st.rerun()
                    else:
                        st.error(result)
                else:
                    st.error("System not initialized. Please set up API keys first.")
    
    with tab2:
        with st.container():
            st.markdown("### Register New Account")
            new_username = st.text_input("Username", key="reg_username")
            new_password = st.text_input("Password", type="password", key="reg_password")
            confirm_password = st.text_input("Confirm Password", type="password", key="confirm_password")
            email = st.text_input("Email")
            
            col1, col2, col3 = st.columns([1, 1, 1])
            with col2:
                register_button = st.button("Register", use_container_width=True)
            
            if register_button:
                if st.session_state.initialized and 'users_collection' in st.session_state:
                    if not new_username or not new_password or not email:
                        st.error("All fields are required")
                    elif new_password != confirm_password:
                        st.error("Passwords don't match")
                    else:
                        with st.spinner("Creating your account..."):
                            success, result = register_user(st.session_state.users_collection, new_username, new_password, email)
                        
                        if success:
                            st.success("Registration successful! Please login with your new credentials.")
                            time.sleep(1)
                            st.rerun()
                        else:
                            st.error(result)
                else:
                    st.error("System not initialized. Please set up API keys first.")

# Main app UI for logged-in users
def show_main_app():
    st.title("🧪 Chemistry Q&A Personal Knowledge Base")
    
    # User info bar at the top
    st.markdown(f"""
    <div class="user-info">
        <strong>Welcome, {st.session_state.username}!</strong> | Your personal chemistry knowledge base
    </div>
    """, unsafe_allow_html=True)
    
    # Initialize user collection if not already done
    if 'user_collection' not in st.session_state:
        try:
            st.session_state.user_collection = get_user_collection(
                token=st.session_state.astra_token,
                api_endpoint=st.session_state.astra_api_endpoint,
                user_id=st.session_state.user_id,
                keyspace=st.session_state.keyspace if 'keyspace' in st.session_state else None,
                dimension=768  # Dimension for the 'all-mpnet-base-v2' model
            )
        except Exception as e:
            st.error(f"Error connecting to your personal database: {e}")
            return
    
    # Check if collection is empty
    collection_empty = check_collection_empty(st.session_state.user_collection)
    
    # Sidebar for logout and file upload
    with st.sidebar:
        # Logout button
        if st.button("Logout"):
            # Clear user session data
            st.session_state.user_logged_in = False
            st.session_state.user_id = None
            st.session_state.username = None
            st.session_state.jwt_token = None
            if 'user_collection' in st.session_state:
                del st.session_state.user_collection
            st.session_state.file_processed = False
            st.session_state.history = []
            st.success("Logged out successfully!")
            time.sleep(1)
            st.rerun()
        
        st.markdown("---")
        st.header("Upload Document")
        st.write("Add chemistry content to your personal knowledge base.")
        
        uploaded_file = st.file_uploader("Choose a PDF or PowerPoint file", type=["pdf", "pptx"])
        
        if uploaded_file is not None:
            if st.button("Process Document"):
                with st.spinner("Processing document..."):
                    questions = process_uploaded_file(uploaded_file)
                    if questions:
                        # Store in user's collection
                        store_qa_in_db(questions, st.session_state.user_collection)
                        st.success(f"Successfully processed document and extracted {len(questions)} questions!")
                        collection_empty = False  # Update status after processing
                        st.rerun()
                    else:
                        st.error(result)
            else:
                st.error("System not initialized. Please set up API keys first.")
    
    # Main area - different view based on whether collection is empty
    if collection_empty:
        st.info("👋 Your personal knowledge base is empty.")
        st.markdown("""
        ### Get Started
        
        To begin using your personal chemistry Q&A system:
        
        1. Upload at least one PDF or PowerPoint file containing chemistry content
        2. The system will extract questions, answers, and explanations
        3. Then you can start asking questions related to the content
        
        **Use the file uploader in the sidebar to add your first document.**
        """)
    else:
        # Two column layout for Q&A area
        col1, col2 = st.columns([3, 1])
        
        with col1:
            # Q&A Interface
            st.header("Ask Chemistry Questions")
            st.write("Ask questions about the chemistry content in your personal knowledge base.")
            
            # Query input
            user_query = st.text_input("Ask a chemistry question:", placeholder="What is the mechanism of the Grignard reaction?")
            
            if st.button("Submit Question"):
                if user_query:
                    # Run RAG pipeline with user's collection
                    answer = rag_pipeline(user_query, st.session_state.user_collection)
                    
                    # Display the answer
                    st.markdown("### Your Question")
                    st.info(user_query)
                    st.markdown("### Answer")
                    st.markdown(answer)
                else:
                    st.warning("Please enter a question first.")
        
        with col2:
            # Recent questions
            if st.session_state.history:
                st.header("Recent Questions")
                for i, item in enumerate(reversed(st.session_state.history[:5])):  # Show last 5 questions
                    with st.expander(f"{item['query'][:30]}..."):
                        st.write(f"**Q:** {item['query']}")
                        st.write(f"**A:** {item['answer'][:100]}...")
        
        # Display full history with expandable sections
        if st.session_state.history:
            st.markdown("---")
            st.header("Question History")
            
            for i, item in enumerate(reversed(st.session_state.history)):
                with st.expander(f"Q: {item['query']}"):
                    st.markdown("#### Question")
                    st.info(item["query"])
                    st.markdown("#### Answer")
                    st.markdown(item["answer"])

# Main Streamlit App
def main():
    # API Keys setup first
    if not st.session_state.initialized:
        st.title("🧪 Chemistry Q&A System Setup")
        
        # Sidebar for setup and configuration
        with st.sidebar:
            st.header("API Configuration")
            st.write("Configuration is now read from environment variables.") # Inform user
            
            if st.button("Initialize System"):
                with st.spinner("Initializing system..."):
                    if initialize_models(): # Call without arguments
                        st.success("System initialized successfully!")
                        time.sleep(1)
                        st.rerun()
                    else:
                        st.error("Failed to initialize. Please check your API credentials in your .env file.") # Updated error message
        
        # Main area instruction
        st.info("👈 Please ensure your API credentials are set in a `.env` file in the same directory.") # Updated info message
        st.markdown("""
        ### Welcome to the Chemistry Q&A System with User Authentication
        
        This application allows:
        1. Users to register and log in to their personal accounts
        2. Each user to build their own chemistry knowledge base
        3. Uploading PDF and PowerPoint files containing chemistry content
        4. Asking questions about the chemistry concepts in their documents
        
        **To get started, please ensure your API credentials are set in a `.env` file and then click "Initialize System" in the sidebar.**
        """)
    
    # After initialization, show either auth page or main app
    elif not st.session_state.user_logged_in:
        show_auth_page()
    else:
        show_main_app()

if __name__ == "__main__":
    main()