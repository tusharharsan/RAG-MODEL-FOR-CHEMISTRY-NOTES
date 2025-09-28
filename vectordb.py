import streamlit as st
import os
import time
from pptx import Presentation
import pytesseract
from PIL import Image
import io
import google.generativeai as genai
import json
import PyPDF2
from sentence_transformers import SentenceTransformer
import re
from astrapy import DataAPIClient
from astrapy.info import CollectionDescriptor
from langchain_astradb import CollectionVectorServiceOptions
# Set page configuration
st.set_page_config(
    page_title="Chemistry Q&A RAG System",
    page_icon="🧪",
    layout="wide"
)

# Apply custom styling
st.markdown("""
<style>
    .main {
        background-color: #000000;
    }
    .stApp {
        max-width: 1200px;
        margin: 0 auto;
    }
    h1 {
        color: #1e3d59;
    }
    .stButton button {
        background-color: #1e3d59;
        color: white;
    }
    .status {
        padding: 10px;
        border-radius: 5px;
        margin-bottom: 10px;
    }
    .success {
        background-color: #d4edda;
        color: #155724;
    }
    .info {
        background-color: #d1ecf1;
        color: #0c5460;
    }
    .results {
        padding: 20px;
        border-radius: 5px;
        background-color: white;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        margin-top: 20px;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state variables if they don't exist
if 'initialized' not in st.session_state:
    st.session_state.initialized = False
    st.session_state.api_key_set = False
    st.session_state.astra_credentials_set = False
    st.session_state.file_processed = False
    st.session_state.collection = None
    st.session_state.embedding_model = None
    st.session_state.qa_count = 0
    st.session_state.questions = []
    st.session_state.history = []

# Astra DB functions
def get_vectorize_providers(token: str, environment: str = None, api_endpoint: str = None):
    """
    Get the list of embedding providers available in Astra DB.
    
    Parameters:
    token (str): Authentication token for accessing Astra DB
    environment (str, optional): The environment for the Astra DB API Endpoint
    api_endpoint (str, optional): The API Endpoint for the Astra DB instance
    
    Returns:
    dict: A mapping of provider display names to their keys and models
    """
    from collections import defaultdict
    from astrapy import AstraDBAdmin
    
    try:
        # Get the admin object
        admin = AstraDBAdmin(token=token, environment=environment)
        db_admin = admin.get_database_admin(api_endpoint=api_endpoint)

        # Get the list of embedding providers
        embedding_providers = db_admin.find_embedding_providers().as_dict()

        vectorize_providers_mapping = {}
        # Map the provider display name to the provider key and models
        for provider_key, provider_data in embedding_providers["embeddingProviders"].items():
            # Get the provider display name and models
            display_name = provider_data["displayName"]
            models = [model["name"] for model in provider_data["models"]]

            # Build our mapping
            vectorize_providers_mapping[display_name] = [provider_key, models]

        # Sort the resulting dictionary
        return defaultdict(list, dict(sorted(vectorize_providers_mapping.items())))
    except Exception as e:
        msg = f"Error fetching vectorize providers: {e}"
        raise ValueError(msg) from e

def get_or_create_user_collection(
    token: str,
    api_endpoint: str,
    user_id: str,
    environment: str = None,
    keyspace: str = None,
    dimension: int = 1536,
    embedding_generation_provider: str = None,
    embedding_generation_model: str = None
):
    """
    Gets an existing collection for a user or creates one if it doesn't exist.
    
    Parameters:
    token (str): Authentication token for accessing Astra DB
    api_endpoint (str): The API Endpoint for the Astra DB instance
    user_id (str): Unique user identifier to append to collection name
    environment (str, optional): The environment for the Astra DB API Endpoint
    keyspace (str, optional): Optional keyspace within Astra DB to use for the collection
    dimension (int, optional): Dimensions of the embeddings (used only for "Bring your own" embedding provider)
    embedding_generation_provider (str, optional): Provider to use for generating embeddings
    embedding_generation_model (str, optional): Model to use for generating embeddings
    
    Returns:
    Collection: The existing or newly created Astra DB collection
    """
    from astrapy import DataAPIClient
    
    # Create the collection name with user ID prefix/suffix
    collection_name = f"user_{user_id}_collection"
    
    # Create the data API client
    client = DataAPIClient(token=token, environment=environment)
    
    # Get the database object
    database = client.get_database(api_endpoint=api_endpoint, token=token, keyspace=keyspace)
    
    # Check if collection exists
    existing_collections = list(database.list_collection_names(keyspace=keyspace))
    
    if collection_name in existing_collections:
        # Collection exists, return it
        print(f"Using existing collection: {collection_name}")
        return database.get_collection(collection_name, keyspace=keyspace)
    
    # Collection doesn't exist, need to create it
    print(f"Creating new collection: {collection_name}")
    
    # Build vectorize options, if needed
    vectorize_options = None
    if embedding_generation_provider and embedding_generation_model:
        # Get vectorize providers mapping
        vectorize_providers = get_vectorize_providers(
            token=token, 
            environment=environment, 
            api_endpoint=api_endpoint
        )
        
        provider_key = vectorize_providers.get(embedding_generation_provider, [None, []])[0]
        
        vectorize_options = CollectionVectorServiceOptions(
            provider=provider_key,
            model_name=embedding_generation_model,
        )
    
    # Create the collection
    return database.create_collection(
        name=collection_name,
        keyspace=keyspace,
        dimension=dimension if not vectorize_options else None,
        service=vectorize_options,
    )

# Initialize models and database
def initialize_models(gemini_api_key, astra_token, astra_api_endpoint, keyspace=None):
    try:
        # Configure Gemini
        genai.configure(api_key=gemini_api_key)
        st.session_state.model = genai.GenerativeModel('gemini-2.0-flash')
        
        # Configure Tesseract (if needed)
        if os.name == 'nt':  # Windows
            pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract'
        
        # Initialize Embedding Model
        st.session_state.embedding_model = SentenceTransformer('all-mpnet-base-v2')
        
        # Initialize Astra DB Collection
        user_id = "chemistry_app"  # A fixed user ID for this app
        st.session_state.collection = get_or_create_user_collection(
            token=astra_token,
            api_endpoint=astra_api_endpoint,
            user_id=user_id,
            keyspace=keyspace,
            dimension=768,  # Dimension for 'all-mpnet-base-v2' model
        )
        
        st.session_state.initialized = True
        st.session_state.api_key_set = True
        st.session_state.astra_credentials_set = True
        return True
    except Exception as e:
        st.error(f"Error initializing models and database: {e}")
        return False

# Functions from the original chemistry app code
def extract_page_data(page, page_type):
    """Extracts text and image captions from a page (PPT or PDF)."""
    if page_type == "ppt":
        slide_data = {"slide_text": [], "image_captions": []}
        for shape in page.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_data["slide_text"].append(paragraph.text)
            elif shape.has_picture:
                try:
                    image_bytes = shape.image.blob
                    image = Image.open(io.BytesIO(image_bytes))
                    text = pytesseract.image_to_string(image)
                    slide_data["image_captions"].append(text)
                except Exception as e_image:
                    st.warning(f"Error processing image: {e_image}")
        return slide_data

    elif page_type == "pdf":
        page_text = page.extract_text()
        return {"slide_text": [page_text], "image_captions": []}  # PDFs don't have images in the same way.

    else:
        return {"slide_text": [], "image_captions": []}

def process_ppt_pages(ppt_file):
    """Processes each page of a PPT and extracts data."""
    try:
        presentation = Presentation(ppt_file)
        page_data = []
        for i, slide in enumerate(presentation.slides):
            data = extract_page_data(slide, "ppt")
            page_data.append({"page_number": i + 1, "data": data, "page_type": "ppt"})
        return page_data
    except Exception as e:
        st.error(f"Error processing PPT: {e}")
        return []

def process_pdf_pages(pdf_file):
    """Processes each page of a PDF and extracts data."""
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        page_data = []
        for i, page in enumerate(pdf_reader.pages):
            data = extract_page_data(page, "pdf")
            page_data.append({"page_number": i + 1, "data": data, "page_type": "pdf"})
        return page_data
    except Exception as e:
        st.error(f"Error processing PDF: {e}")
        return []

def extract_chemistry_qa_from_page(page_data):
    """Extracts chemistry questions and answers from a page's text using Gemini."""
    slide_text = "\n".join(page_data["slide_text"])
    image_captions = "\n".join(page_data["image_captions"])
    combined_text = slide_text + "\n" + image_captions

    if not combined_text.strip():
        return None

    prompt = f"""
    Extract chemistry questions and answers from the following text. 
    Format your response as a JSON object with the following structure:
    {{
        "questions": [
            {{
                "question_text": "Full text of the question",
                "options": {{
                    "A": "Text of option A",
                    "B": "Text of option B",
                    "C": "Text of option C",
                    "D": "Text of option D"
                }},
                "correct_answer": "Letter of the correct option (A, B, C, or D)",
                "explanation": "Explanation of the answer"
            }},
            ...
        ]
    }}
    
    If there are no clear questions and answers, return an empty questions array.
    If options are not lettered, assign them letters in order.
    
    Text: {combined_text}
    """

    try:
        response = st.session_state.model.generate_content(prompt)
        response_text = response.text
        
        # Try to parse as JSON
        try:
            # Find JSON content within response (in case there's extra text)
            json_match = re.search(r'```json\s*([\s\S]*?)\s*```', response_text)
            if json_match:
                json_str = json_match.group(1)
            else:
                json_str = response_text
                
            # Clean up any markdown formatting
            json_str = re.sub(r'```.*?```', '', json_str, flags=re.DOTALL)
            
            # Parse the JSON
            result = json.loads(json_str)
            return result
        except json.JSONDecodeError:
            # If parsing fails, try to identify and fix common JSON issues
            cleaned_text = response_text.strip()
            # Remove markdown code blocks if present
            cleaned_text = re.sub(r'```json|```', '', cleaned_text)
            
            # Attempt to parse again
            try:
                result = json.loads(cleaned_text)
                return result
            except:
                st.warning(f"Failed to parse JSON from response: {response_text}")
                # Return a structured empty result
                return {"questions": []}
    except Exception as e:
        st.error(f"Error extracting Q&A: {e}")
        return {"questions": []}

def process_uploaded_file(uploaded_file):
    """Process the uploaded file and extract Q&A data."""
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    
    # Create a progress bar
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    if file_extension == '.pptx':
        # Save the uploaded file temporarily
        with open("temp_file.pptx", "wb") as f:
            f.write(uploaded_file.getbuffer())
        status_text.text("Processing PowerPoint file...")
        page_data = process_ppt_pages("temp_file.pptx")
        os.remove("temp_file.pptx")  # Clean up
    elif file_extension == '.pdf':
        status_text.text("Processing PDF file...")
        page_data = process_pdf_pages(uploaded_file)
    else:
        st.error("Unsupported file type. Please upload a PDF or PPTX file.")
        return []
    
    all_questions = []
    total_pages = len(page_data)
    
    for i, page in enumerate(page_data):
        status_text.text(f"Extracting Q&A from page {i+1}/{total_pages}...")
        progress_bar.progress((i + 0.5) / total_pages)
        
        qa_data = extract_chemistry_qa_from_page(page["data"])
        
        if qa_data and "questions" in qa_data and qa_data["questions"]:
            for q in qa_data["questions"]:
                q["page_number"] = page["page_number"]
                q["page_type"] = page["page_type"]
                all_questions.append(q)
        
        progress_bar.progress((i + 1) / total_pages)
    
    status_text.text(f"Extracted {len(all_questions)} questions total.")
    time.sleep(1)  # Give users time to see the message
    status_text.empty()
    progress_bar.empty()
    
    return all_questions

def store_qa_in_db(questions):
    """Stores the Q&A data in Astra DB."""
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    # Clear collection by deleting all documents (if any)
    try:
        # Use delete_many to remove all documents in the collection
        st.session_state.collection.delete_many({})
        st.success("Collection cleared successfully.")
    except Exception as e:
        st.warning(f"Could not clear existing collection: {e}")
    
    status_text.text("Storing questions in Astra DB...")
    
    for i, q in enumerate(questions):
        # Update progress
        progress_bar.progress((i + 0.5) / len(questions))
        
        # Prepare data for embedding
        question_id = f"q{i+1}"
        
        # Create full text representation for embedding
        full_text = f"Question: {q['question_text']}\n"
        
        if "options" in q:
            options_text = ""
            for option_key, option_text in q["options"].items():
                options_text += f"Option {option_key}: {option_text}\n"
            full_text += options_text
        
        if "correct_answer" in q:
            full_text += f"Correct Answer: {q['correct_answer']}\n"
        
        if "explanation" in q:
            full_text += f"Explanation: {q['explanation']}\n"
        
        # Generate embeddings
        embedding = st.session_state.embedding_model.encode(full_text).tolist()
        
        # Store in Astra DB
        document = {
            "text": full_text,
            "page_number": q.get("page_number", 0),
            "page_type": q.get("page_type", "unknown"),
            "question_id": question_id,
            "$vector": embedding
        }
        
        st.session_state.collection.insert_one(document)
        progress_bar.progress((i + 1) / len(questions))
    
    status_text.text(f"Successfully stored {len(questions)} questions in Astra DB.")
    time.sleep(1)  # Give users time to see the message
    status_text.empty()
    progress_bar.empty()
    
    st.session_state.qa_count = len(questions)
    st.session_state.questions = questions
    st.session_state.file_processed = True

def query_qa_db(query, n_results=5):
    """Queries the chemistry Q&A vector database and returns relevant results."""
    # Generate embedding for the query
    query_embedding = st.session_state.embedding_model.encode(query).tolist()
    
    # Search in Astra DB
    results = st.session_state.collection.find(
        vector=query_embedding,
        limit=n_results
    )
    
    # Convert cursor to list of documents
    documents = list(results)
    
    # Format the results to match what the generate_answer function expects
    formatted_results = {
        "documents": [[doc.get("text", "")] for doc in documents],
        "metadatas": documents
    }
    
    return formatted_results

def generate_answer(user_query, retrieved_results):
    """Generates an answer using Gemini based on retrieved results."""
    # Compile context from retrieved documents
    context_docs = []
    for doc_list in retrieved_results.get("documents", []):
        if doc_list:  # Make sure it's not empty
            context_docs.extend(doc_list)
    
    context = "\n\n".join(context_docs)
    
    prompt = f"""
    Based on the following chemistry question and answer information, provide a detailed answer to the user's query.
    
    Retrieved Information:
    {context}
    
    User Query: {user_query}
    
    Provide a clear, detailed answer with explanations of the chemistry concepts involved. If the answer involves chemical reactions, explain the mechanism where appropriate.
    """
    
    response = st.session_state.model.generate_content(prompt)
    return response.text

def rag_pipeline(user_query, n_results=3):
    """Full RAG pipeline: Retrieval + Answer Generation."""
    # Retrieve relevant Q&A content
    with st.spinner("Searching the knowledge base..."):
        retrieved_results = query_qa_db(user_query, n_results)
    
    # Generate answer based on retrieved content
    with st.spinner("Generating your answer..."):
        answer = generate_answer(user_query, retrieved_results)
    
    # Store in history
    st.session_state.history.append({
        "query": user_query,
        "answer": answer
    })
    
    return answer

def process_text_input(text_content):
    """Process direct text input from the user and extract Q&A data."""
    
    status_text = st.empty()
    progress_bar = st.progress(0)
    
    status_text.text("Processing text input...")
    
    # Create a single page data structure for the text
    page_data = [{
        "page_number": 1,
        "data": {"slide_text": [text_content], "image_captions": []},
        "page_type": "text"
    }]
    
    all_questions = []
    
    status_text.text("Extracting Q&A from text...")
    progress_bar.progress(0.5)
    
    qa_data = extract_chemistry_qa_from_page(page_data[0]["data"])
    
    if qa_data and "questions" in qa_data and qa_data["questions"]:
        for q in qa_data["questions"]:
            q["page_number"] = 1
            q["page_type"] = "text"
            all_questions.append(q)
    
    progress_bar.progress(1.0)
    status_text.text(f"Extracted {len(all_questions)} questions from text input.")
    time.sleep(1)
    status_text.empty()
    progress_bar.empty()
    
    return all_questions

def process_multiple_files(uploaded_files):
    """Process multiple uploaded files and extract Q&A data."""
    all_questions = []
    
    # Create a progress bar for overall progress
    overall_progress = st.progress(0)
    status_text = st.empty()
    
    for i, uploaded_file in enumerate(uploaded_files):
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()
        status_text.text(f"Processing file {i+1}/{len(uploaded_files)}: {uploaded_file.name}")
        
        file_questions = []
        if file_extension == '.pptx':
            # Save the uploaded file temporarily
            with open(f"temp_file_{i}.pptx", "wb") as f:
                f.write(uploaded_file.getbuffer())
            file_questions = process_ppt_pages(f"temp_file_{i}.pptx")
            os.remove(f"temp_file_{i}.pptx")  # Clean up
        elif file_extension == '.pdf':
            file_questions = process_pdf_pages(uploaded_file)
        else:
            st.warning(f"Skipping unsupported file: {uploaded_file.name}. Please upload PDF or PPTX files only.")
            continue
        
        # Process each page in the file
        file_processed_questions = []
        for page in file_questions:
            qa_data = extract_chemistry_qa_from_page(page["data"])
            if qa_data and "questions" in qa_data and qa_data["questions"]:
                for q in qa_data["questions"]:
                    q["page_number"] = page["page_number"]
                    q["page_type"] = page["page_type"]
                    q["source_file"] = uploaded_file.name
                    file_processed_questions.append(q)
        
        all_questions.extend(file_processed_questions)
        status_text.text(f"Extracted {len(file_processed_questions)} questions from {uploaded_file.name}")
        overall_progress.progress((i + 1) / len(uploaded_files))
    
    status_text.text(f"Total questions extracted from all files: {len(all_questions)}")
    time.sleep(1)
    status_text.empty()
    overall_progress.empty()
    
    return all_questions

# Main Streamlit App
def main():
    st.title("🧪 Chemistry Q&A RAG System with Astra DB")
    st.markdown("Upload your chemistry documents and ask questions to get detailed answers!")
    
    # Sidebar for setup and configuration
    with st.sidebar:
        st.header("Setup")
        
        # Gemini API Key input
        gemini_api_key = st.text_input("Enter your Gemini API Key:", type="password", 
                               value=os.getenv("GEMINI_API_KEY", "") if st.session_state.api_key_set else "")
        
        # Astra DB credentials
        st.subheader("Astra DB Configuration")
        astra_token = st.text_input("Astra DB Token:", type="password", 
                                   value="AstraCS:xyz" if st.session_state.astra_credentials_set else "")
        astra_api_endpoint = st.text_input("Astra DB API Endpoint:", 
                                         value="https://your-database-id-region.apps.astra.datastax.com" if st.session_state.astra_credentials_set else "")
        keyspace = st.text_input("Keyspace (optional):", value="")
        
        if gemini_api_key and astra_token and astra_api_endpoint and not st.session_state.initialized:
            if st.button("Initialize System"):
                with st.spinner("Initializing models and database..."):
                    if initialize_models(gemini_api_key, astra_token, astra_api_endpoint, keyspace if keyspace else None):
                        st.success("System initialized successfully!")
                    else:
                        st.error("Failed to initialize. Check your credentials and try again.")
        
        # Only show file upload if system is initialized
        if st.session_state.initialized:
            st.header("Add Content")
            
            # Create tabs for different input methods
            input_tab1, input_tab2 = st.tabs(["File Upload", "Text Input"])
            
            with input_tab1:
                uploaded_files = st.file_uploader("Choose PDF or PowerPoint files", type=["pdf", "pptx"], accept_multiple_files=True)
                
                if uploaded_files:
                    if st.button("Process Files"):
                        with st.spinner("Processing files..."):
                            questions = process_multiple_files(uploaded_files)
                            if questions:
                                store_qa_in_db(questions)
                                st.success(f"Successfully processed {len(uploaded_files)} files and extracted {len(questions)} questions!")
                            else:
                                st.warning("No questions were extracted from the files.")
            
            with input_tab2:
                text_input = st.text_area("Enter chemistry content directly:", height=300, 
                                        placeholder="Paste or type chemistry content here...")
                
                if text_input:
                    if st.button("Process Text"):
                        with st.spinner("Processing text input..."):
                            questions = process_text_input(text_input)
                            if questions:
                                store_qa_in_db(questions)
                                st.success(f"Successfully processed text input and extracted {len(questions)} questions!")
                            else:
                                st.warning("No questions were extracted from the text input.")
        
        # Display stats if file is processed
        if st.session_state.file_processed:
            st.header("Statistics")
            st.info(f"Questions extracted: {st.session_state.qa_count}")
            st.info(f"Questions answered: {len(st.session_state.history)}")
        
        # Clear history button
        if st.session_state.history:
            if st.button("Clear Chat History"):
                st.session_state.history = []
                st.success("Chat history cleared!")
    
    # Main area for Q&A interaction
    if not st.session_state.initialized:
        st.info("👈 Please enter your API key and Astra DB credentials to initialize the system.")
    elif not st.session_state.file_processed:
        st.info("👈 Please upload and process chemistry documents or input text to start asking questions.")
    else:
        # Q&A Interface
        st.header("Ask Chemistry Questions")
        
        # Query input
        user_query = st.text_input("Type your chemistry question here:", placeholder="What is the mechanism of the Grignard reaction?")
        
        col1, col2 = st.columns([1, 5])
        with col1:
            submit_button = st.button("Ask")
        
        # Process the query if button is clicked
        if submit_button and user_query:
            answer = rag_pipeline(user_query)
            
            # Display the answer (latest interaction is shown immediately)
            st.markdown("### Your Question")
            st.info(user_query)
            st.markdown("### Answer")
            st.markdown(answer)
        
        # Display conversation history (excluding the latest if it was just added)
        if len(st.session_state.history) > 1:
            st.markdown("---")
            st.markdown("### Previous Questions")
            
            for i, item in enumerate(st.session_state.history[:-1][::-1]):  # Reverse to show newest first, excluding the latest
                with st.expander(f"Q: {item['query'][:80]}{'...' if len(item['query']) > 80 else ''}"):
                    st.markdown("#### Question")
                    st.info(item["query"])
                    st.markdown("#### Answer")
                    st.markdown(item["answer"])

if __name__ == "__main__":
    main()